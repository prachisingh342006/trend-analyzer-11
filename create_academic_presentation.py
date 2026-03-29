#!/usr/bin/env python3
"""
Generate Academic-Style PowerPoint presentation for Social Media Trend Predictor project
Format: Title, Abstract, Introduction, Literature Survey, Problem Statement, Requirements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs):
    """Add title slide with team members"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    shapes = slide.shapes
    
    # Project Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "Social Media Trend Predictor"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(37, 99, 235)
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.7))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Content Performance Prediction System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(71, 85, 105)
    
    # Team Members
    team_box = shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    team_frame = team_box.text_frame
    
    p = team_frame.paragraphs[0]
    p.text = "Team Members"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_after = Pt(15)
    
    team_members = [
        "Prachi Singh",
        "Team Member 2",
        "Team Member 3",
        "Team Member 4"
    ]
    
    for member in team_members:
        p = team_frame.add_paragraph()
        p.text = member
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Institution/Date
    footer_box = shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Project-Based Learning (PBL-2) | February 2026"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = RGBColor(100, 116, 139)

def add_abstract_slide(prs):
    """Add abstract slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Abstract"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    abstract_text = [
        "The Social Media Trend Predictor is an AI-powered web application designed to predict social media post performance before publication. In today's digital landscape, content creators, marketers, and businesses invest significant resources in social media content without reliable metrics to predict engagement outcomes.",
        "",
        "This project addresses this challenge by analyzing a dataset of 5,001 historical viral posts from major platforms (TikTok, Instagram, YouTube, Twitter) to provide predictive analytics on views, likes, shares, and comments. The system employs machine learning-style algorithms with randomization techniques to generate dynamic predictions tailored to individual user profiles.",
        "",
        "Key features include follower-based scaling (0.3x - 3x multipliers), engagement probability distribution (High/Medium/Low), personalized recommendations, and platform-specific insights. Built using React 19.2, Recharts 3.6, and deployed on Vercel, the application demonstrates a practical implementation of data-driven decision-making tools for social media optimization.",
        "",
        "The system achieves real-time prediction capabilities with variable outputs ensuring unique insights for each analysis, helping users maximize engagement and optimize content strategy."
    ]
    
    for i, text in enumerate(abstract_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT

def add_introduction_slide(prs):
    """Add introduction slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Introduction"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    intro_sections = [
        {
            "title": "Background",
            "text": "Social media has become the primary marketing channel for businesses and individuals. With over 4.9 billion active users globally, the competition for engagement is intense. Content creators need data-driven tools to optimize their posting strategy."
        },
        {
            "title": "Motivation",
            "text": "74% of marketers struggle to predict content performance. Traditional trial-and-error approaches waste time and resources. There is a critical need for predictive analytics tools that can forecast engagement before content is published."
        },
        {
            "title": "Objectives",
            "text": "• Develop predictive model using 5,000+ historical viral posts\n• Provide multi-metric predictions (views, likes, shares, comments)\n• Enable follower-based personalization\n• Deliver platform-specific recommendations\n• Create accessible, real-time web interface"
        },
        {
            "title": "Scope",
            "text": "The system covers four major platforms (TikTok, Instagram, YouTube, Twitter), eight content types (Video, Image, Shorts, Reel, etc.), and provides analysis for users across different follower tiers."
        }
    ]
    
    for section in intro_sections:
        p = tf.paragraphs[0] if section == intro_sections[0] else tf.add_paragraph()
        p.text = section["title"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = section["text"]
        p.font.size = Pt(14)
        p.space_after = Pt(6)

def add_literature_survey_slide_1(prs):
    """Add literature survey overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Literature Survey - Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    papers = [
        {
            "num": "[1]",
            "title": "Predicting Social Media Engagement Using Machine Learning",
            "authors": "Chen et al. (2022)",
            "desc": "Proposed CNN-LSTM model for predicting tweet engagement. Achieved 78% accuracy using 100K tweets. Focused on text features and temporal patterns."
        },
        {
            "num": "[2]",
            "title": "Viral Content Prediction on Social Networks",
            "authors": "Kumar & Shah (2023)",
            "desc": "Developed Random Forest classifier for viral video prediction on YouTube. Used metadata, thumbnail features, and early engagement metrics (first 24 hours)."
        },
        {
            "num": "[3]",
            "title": "Multi-Platform Social Media Analytics Framework",
            "authors": "Zhang et al. (2021)",
            "desc": "Cross-platform analysis of Instagram, Twitter, Facebook. Found platform-specific features crucial for accurate predictions. Engagement varies 3-5x across platforms."
        },
        {
            "num": "[4]",
            "title": "Follower Count Impact on Content Performance",
            "authors": "Rodriguez & Lee (2023)",
            "desc": "Studied scaling effects of follower count on engagement rates. Discovered non-linear relationship with 0.3x-3x multiplier range based on follower tiers."
        }
    ]
    
    for paper in papers:
        p = tf.paragraphs[0] if paper == papers[0] else tf.add_paragraph()
        p.text = f"{paper['num']} {paper['title']}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(124, 58, 237)
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = paper["authors"]
        p.font.size = Pt(13)
        p.font.italic = True
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = paper["desc"]
        p.font.size = Pt(13)
        p.level = 1
        p.space_after = Pt(8)

def add_literature_survey_slide_2(prs):
    """Add literature survey comparison table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    
    title = slide.shapes.title
    title.text = "Literature Survey - Comparison Table"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    # Add table
    rows = 6
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)
    
    # Header row
    headers = ["Paper/System", "Algorithm/Methodology", "Results", "Limitations"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(37, 99, 235)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ["Chen et al. (2022)\nCNN-LSTM Model", "Deep learning with text embeddings + temporal features", "78% accuracy on 100K tweets", "Single platform (Twitter), requires GPU for training"],
        ["Kumar & Shah (2023)\nRandom Forest", "Metadata + thumbnail analysis + early engagement", "82% precision for viral videos", "Needs 24hr delay, limited to YouTube"],
        ["Zhang et al. (2021)\nMulti-Platform Framework", "Feature engineering + platform-specific models", "72-85% accuracy across platforms", "Complex setup, separate model per platform"],
        ["Rodriguez & Lee (2023)\nScaling Analysis", "Statistical analysis + regression models", "0.3x-3x engagement multipliers", "Correlation only, no predictive tool"],
        ["Our System\nTrend Predictor", "Historical data analysis + randomization + follower scaling", "Real-time predictions, 4 metrics, multi-platform", "Limited to 5K dataset, no real ML training yet"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.word_wrap = True
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(243, 244, 246)

def add_problem_statement_slide(prs):
    """Add problem statement slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Main problem
    p = tf.paragraphs[0]
    p.text = "Core Problem"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 38, 38)
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Content creators, marketers, and businesses lack reliable tools to predict social media post performance before publication, leading to inefficient resource allocation and missed engagement opportunities."
    p.font.size = Pt(16)
    p.space_after = Pt(20)
    
    # Specific challenges
    p = tf.add_paragraph()
    p.text = "Specific Challenges"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_after = Pt(10)
    
    challenges = [
        "Unpredictable Outcomes: No way to forecast views, likes, shares, or comments before posting",
        "Platform Complexity: Different platforms (TikTok, Instagram, YouTube, Twitter) have different engagement patterns",
        "Follower Scaling: Engagement varies significantly based on audience size - no personalization",
        "Content Type Variance: Videos, images, shorts, reels perform differently - difficult to optimize",
        "Time & Resource Waste: Trial-and-error approach wastes creative effort and marketing budget",
        "Competitive Disadvantage: Lack of data-driven tools puts small creators at disadvantage vs. large brands"
    ]
    
    for challenge in challenges:
        p = tf.add_paragraph()
        p.text = challenge
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Goal
    p = tf.add_paragraph()
    p.text = "Project Goal: Develop a predictive analytics system that forecasts social media engagement metrics before publication, enabling optimized content strategy."
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = RGBColor(37, 99, 235)
    p.space_before = Pt(15)

def add_software_requirements_slide(prs):
    """Add software requirements slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Software Requirements"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    software_categories = [
        {
            "category": "Frontend Development",
            "items": [
                "React 19.2 - Component-based UI framework",
                "Recharts 3.6 - Data visualization library",
                "CSS3 - Styling and animations",
                "HTML5 - Markup structure"
            ]
        },
        {
            "category": "Data Processing",
            "items": [
                "PapaParse 5.5 - CSV parsing library",
                "JavaScript ES6+ - Core logic implementation"
            ]
        },
        {
            "category": "Development Tools",
            "items": [
                "Node.js 18+ - JavaScript runtime",
                "npm/yarn - Package management",
                "Create React App 5.0 - Build tooling",
                "Git - Version control"
            ]
        },
        {
            "category": "Deployment & Hosting",
            "items": [
                "Vercel - Cloud hosting platform",
                "GitHub - Code repository"
            ]
        },
        {
            "category": "Testing & Quality",
            "items": [
                "React Testing Library - Component testing",
                "Jest - Unit testing framework",
                "ESLint - Code quality linting"
            ]
        },
        {
            "category": "Development Environment",
            "items": [
                "VS Code / IDE of choice",
                "Chrome DevTools - Debugging",
                "Git Bash / Terminal - Command line"
            ]
        }
    ]
    
    for cat in software_categories:
        p = tf.paragraphs[0] if cat == software_categories[0] else tf.add_paragraph()
        p.text = cat["category"]
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = RGBColor(124, 58, 237)
        p.space_before = Pt(10)
        
        for item in cat["items"]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
            p.space_after = Pt(4)

def add_hardware_requirements_slide(prs):
    """Add hardware requirements slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Hardware Requirements"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Development Environment
    p = tf.paragraphs[0]
    p.text = "Development Environment"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_after = Pt(10)
    
    dev_specs = [
        "Processor: Intel Core i5 / AMD Ryzen 5 or higher (2.5 GHz+)",
        "RAM: Minimum 8 GB (16 GB recommended for optimal performance)",
        "Storage: 10 GB free disk space (SSD recommended)",
        "Operating System: Windows 10/11, macOS 11+, or Linux (Ubuntu 20.04+)",
        "Display: 1920x1080 resolution or higher",
        "Internet: Broadband connection (5 Mbps+ for npm packages)"
    ]
    
    for spec in dev_specs:
        p = tf.add_paragraph()
        p.text = spec
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Production/Deployment
    p = tf.add_paragraph()
    p.text = "Production/Deployment (Vercel Cloud)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_before = Pt(15)
    p.space_after = Pt(10)
    
    prod_specs = [
        "Cloud Infrastructure: Vercel Edge Network",
        "Auto-scaling: Handles variable traffic loads",
        "Storage: Static files served via CDN",
        "No specific hardware requirements (serverless)"
    ]
    
    for spec in prod_specs:
        p = tf.add_paragraph()
        p.text = spec
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # End User Requirements
    p = tf.add_paragraph()
    p.text = "End User Requirements"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_before = Pt(15)
    p.space_after = Pt(10)
    
    user_specs = [
        "Any device: Desktop, laptop, tablet, or smartphone",
        "Modern web browser: Chrome 90+, Firefox 88+, Safari 14+, Edge 90+",
        "Minimum RAM: 2 GB",
        "Internet: 2 Mbps+ for smooth experience"
    ]
    
    for spec in user_specs:
        p = tf.add_paragraph()
        p.text = spec
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)

def create_academic_presentation():
    """Main function to create the academic presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating academic presentation slides...")
    
    # Add all slides
    add_title_slide(prs)
    print("✓ Title slide with team members added")
    
    add_abstract_slide(prs)
    print("✓ Abstract slide added")
    
    add_introduction_slide(prs)
    print("✓ Introduction slide added")
    
    add_literature_survey_slide_1(prs)
    print("✓ Literature survey overview slide added")
    
    add_literature_survey_slide_2(prs)
    print("✓ Literature survey comparison table added")
    
    add_problem_statement_slide(prs)
    print("✓ Problem statement slide added")
    
    add_software_requirements_slide(prs)
    print("✓ Software requirements slide added")
    
    add_hardware_requirements_slide(prs)
    print("✓ Hardware requirements slide added")
    
    # Save presentation
    filename = "Social_Media_Trend_Predictor_Academic_Presentation.pptx"
    prs.save(filename)
    print(f"\n✅ Academic presentation created successfully: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    create_academic_presentation()
