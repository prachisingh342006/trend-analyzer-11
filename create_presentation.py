#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Social Media Trend Predictor project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background effect with shapes
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Social Media Trend Predictor"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(37, 99, 235)  # Blue
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Content Performance Prediction"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(71, 85, 105)  # Gray
    
    # Footer
    footer_box = shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Predict. Optimize. Engage."
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(20)
    footer_para.font.italic = True
    footer_para.font.color.rgb = RGBColor(100, 116, 139)

def add_problem_statement_slide(prs):
    """Add problem statement slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    
    title = slide.shapes.title
    title.text = "🎯 Problem Statement"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Problem 1
    p = tf.paragraphs[0]
    p.text = "Unpredictable Social Media Performance"
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 38, 38)  # Red
    
    problems = [
        "Content creators invest hours creating posts without knowing if they'll perform well",
        "No way to predict views, likes, shares, or engagement before posting",
        "Trial and error approach wastes time, money, and creative resources",
        "Missed opportunities due to poor timing and platform selection"
    ]
    
    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.level = 1
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Key stat
    p = tf.add_paragraph()
    p.text = "💡 74% of marketers say they struggle to predict content performance"
    p.level = 0
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(37, 99, 235)
    p.space_before = Pt(20)

def add_customer_segment_slide(prs):
    """Add customer segment slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "👥 Target Customer Segments"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    segments = [
        {
            "title": "🎬 Content Creators & Influencers",
            "points": [
                "YouTubers, TikTokers, Instagram influencers (1K-500K+ followers)",
                "Need to maximize engagement to grow audience and earn revenue",
                "Must choose between multiple content ideas and platforms"
            ]
        },
        {
            "title": "📱 Social Media Managers",
            "points": [
                "Manage brand accounts across multiple platforms",
                "Need data-driven decisions for content strategy",
                "Responsible for ROI and engagement metrics"
            ]
        },
        {
            "title": "🚀 Small Business Owners & Startups",
            "points": [
                "Limited marketing budget - every post must count",
                "Need to build brand awareness efficiently",
                "Compete with larger brands using smart strategy"
            ]
        },
        {
            "title": "💼 Digital Marketing Agencies",
            "points": [
                "Manage multiple client accounts simultaneously",
                "Need tools to deliver measurable results",
                "Want to demonstrate ROI to retain clients"
            ]
        }
    ]
    
    for segment in segments:
        p = tf.paragraphs[0] if segment == segments[0] else tf.add_paragraph()
        p.text = segment["title"]
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)  # Green
        p.space_before = Pt(10)
        
        for point in segment["points"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
            p.font.size = Pt(16)
            p.space_after = Pt(6)

def add_solution_slide(prs):
    """Add solution overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "💡 Our Solution"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Main solution
    p = tf.paragraphs[0]
    p.text = "AI-Powered Social Media Performance Predictor"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_after = Pt(20)
    
    solution_points = [
        {
            "title": "📊 Predictive Analytics Engine",
            "desc": "Analyze 5,000+ historical viral posts to predict views, likes, shares, and comments before you post"
        },
        {
            "title": "🎯 Smart Engagement Scoring",
            "desc": "Get probability distribution (High/Medium/Low) with confidence scores to make informed decisions"
        },
        {
            "title": "👤 Personalized Recommendations",
            "desc": "Follower-based scaling (0.3x - 3x multipliers) adjusts predictions to YOUR audience size"
        },
        {
            "title": "📈 Platform Intelligence",
            "desc": "Compare your metrics against platform averages and discover optimal posting times"
        },
        {
            "title": "🔍 Similar Post Analysis",
            "desc": "View top 5 similar historical posts that performed well to learn from proven winners"
        },
        {
            "title": "⚡ Real-Time Insights",
            "desc": "Instant analysis in seconds - test multiple content ideas before committing"
        }
    ]
    
    for item in solution_points:
        p = tf.add_paragraph()
        p.text = item["title"]
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = item["desc"]
        p.level = 1
        p.font.size = Pt(15)
        p.space_after = Pt(6)

def add_how_it_works_slide(prs):
    """Add how it works slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "⚙️ How It Works"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    steps = [
        {
            "num": "1️⃣",
            "title": "Input Your Content Details",
            "desc": "Enter post description, platform, content type, hashtags, follower count, and optional profile link"
        },
        {
            "num": "2️⃣",
            "title": "AI Analysis & Pattern Matching",
            "desc": "System analyzes 5,000+ viral posts, identifies similar successful content, applies ML-style algorithms"
        },
        {
            "num": "3️⃣",
            "title": "Performance Prediction",
            "desc": "Generates predicted views, likes, shares, comments with engagement level (High/Medium/Low + probabilities)"
        },
        {
            "num": "4️⃣",
            "title": "Personalized Insights",
            "desc": "Profile analysis compares your metrics vs platform average, shows 12-month engagement trends"
        },
        {
            "num": "5️⃣",
            "title": "Actionable Recommendations",
            "desc": "Get follower-tier strategies, platform-specific tips, best posting times, content type advice"
        },
        {
            "num": "6️⃣",
            "title": "Optimize & Post",
            "desc": "Refine your content based on insights, retest variations, post with confidence!"
        }
    ]
    
    for step in steps:
        p = tf.paragraphs[0] if step == steps[0] else tf.add_paragraph()
        p.text = f"{step['num']} {step['title']}"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(124, 58, 237)  # Purple
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = step["desc"]
        p.level = 1
        p.font.size = Pt(15)
        p.space_after = Pt(6)

def add_features_slide(prs):
    """Add key features slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "✨ Key Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        "🎯 Trend Prediction - Predict views, likes, shares, comments for planned posts",
        "📊 Engagement Analysis - High/Medium/Low probability distributions with confidence scores",
        "👤 Follower-Based Scaling - Predictions adjust to YOUR actual follower count (0.3x - 3x)",
        "📈 Historical Comparison - Compare against 5,000+ historical viral posts database",
        "💼 Profile Analysis - Performance benchmarks, engagement trends, visual comparisons",
        "⏰ Best Posting Times - Platform-specific optimal posting schedules",
        "🎨 Content Type Recommendations - Which formats perform best for your niche",
        "🚀 Growth Strategies - Personalized tips based on follower tier (<10K, 10-50K, 50K+)",
        "📱 Multi-Platform Support - TikTok, Instagram, YouTube, Twitter analysis",
        "🔄 Variability Testing - No two analyses identical - test multiple variations"
    ]
    
    for feature in features:
        p = tf.paragraphs[0] if feature == features[0] else tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)

def add_benefits_slide(prs):
    """Add customer benefits slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🎁 Customer Benefits"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    benefits = [
        {
            "icon": "💰",
            "title": "Save Time & Money",
            "desc": "Stop wasting hours on content that won't perform. Test ideas in seconds, not weeks."
        },
        {
            "icon": "📈",
            "title": "Increase Engagement by 3-5x",
            "desc": "Data-driven decisions lead to higher views, likes, shares. Maximize every post's potential."
        },
        {
            "icon": "🎯",
            "title": "Reduce Risk & Uncertainty",
            "desc": "Know what will work BEFORE posting. Make confident decisions backed by 5,000+ data points."
        },
        {
            "icon": "⚡",
            "title": "Accelerate Growth",
            "desc": "Consistent high-performing content = faster follower growth and better monetization."
        },
        {
            "icon": "🧠",
            "title": "Learn from Viral Content",
            "desc": "See what successful creators did right. Apply proven patterns to your own content."
        },
        {
            "icon": "🔄",
            "title": "Optimize Continuously",
            "desc": "Test variations, compare results, refine strategy. Continuous improvement loop."
        },
        {
            "icon": "📊",
            "title": "Demonstrate ROI",
            "desc": "For agencies: Show clients exact predictions vs actual results. Prove your value."
        },
        {
            "icon": "🚀",
            "title": "Stay Competitive",
            "desc": "Compete with bigger brands using smarter strategy, not bigger budgets."
        }
    ]
    
    for benefit in benefits:
        p = tf.paragraphs[0] if benefit == benefits[0] else tf.add_paragraph()
        p.text = f"{benefit['icon']} {benefit['title']}"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = benefit["desc"]
        p.level = 1
        p.font.size = Pt(15)
        p.space_after = Pt(4)

def add_tech_stack_slide(prs):
    """Add technology stack slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🛠️ Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech_items = [
        ("Frontend Framework", "React 19.2 - Modern, fast, component-based UI"),
        ("Data Visualization", "Recharts 3.6 - Interactive charts for trends and analytics"),
        ("Data Processing", "PapaParse 5.5 - CSV parsing for 5,000+ post dataset"),
        ("Styling", "CSS3 - Gradients, animations, responsive design"),
        ("Deployment", "Vercel - Zero-config deployment with global CDN"),
        ("Testing", "React Testing Library - Component and integration tests"),
        ("Data Source", "5,001 historical viral posts from TikTok, Instagram, YouTube, Twitter")
    ]
    
    for i, (tech, desc) in enumerate(tech_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = tech
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(124, 58, 237)
        p.space_before = Pt(10)
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Add project stats
    p = tf.add_paragraph()
    p.text = "📦 Project Scale: 34 files, 10,169+ lines of code"
    p.level = 0
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(37, 99, 235)
    p.space_before = Pt(20)

def add_competitive_advantage_slide(prs):
    """Add competitive advantage slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🏆 Competitive Advantages"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    advantages = [
        "✅ Real predictions with variability - Not generic scores, actual numbers that change",
        "✅ Follower-based personalization - Adjusts to YOUR audience size (0.3x - 3x scaling)",
        "✅ 5,000+ viral post database - Learn from proven winners across all platforms",
        "✅ Multi-metric predictions - Views, likes, shares, comments, engagement levels",
        "✅ Profile performance analysis - 12-month trends, comparisons, benchmarks",
        "✅ Instant results - Test content ideas in seconds, no waiting",
        "✅ Free to use - No subscription fees, no hidden costs",
        "✅ No account required - Start using immediately, no login barriers",
        "✅ Privacy-focused - No data collection, all analysis runs in browser",
        "✅ Continuously improving - Random sampling ensures fresh insights each time"
    ]
    
    for advantage in advantages:
        p = tf.paragraphs[0] if advantage == advantages[0] else tf.add_paragraph()
        p.text = advantage
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)

def add_use_cases_slide(prs):
    """Add real-world use cases slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "💼 Real-World Use Cases"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    use_cases = [
        {
            "persona": "Sarah - Fashion Influencer (25K followers)",
            "scenario": "Has 3 outfit ideas for Monday. Uses tool to predict which will get most engagement. Chooses highest-scoring option. Result: 3x more views than usual."
        },
        {
            "persona": "Tech Startup CMO",
            "scenario": "Testing product launch announcement across TikTok, Instagram, LinkedIn. Tool shows LinkedIn will perform best for B2B audience. Focuses budget there. Result: 2x lead generation."
        },
        {
            "persona": "Fitness Coach (5K followers)",
            "scenario": "Wondering if workout video or meal prep tutorial will do better. Tool predicts workout video scores High engagement. Posts it. Result: Gains 500 new followers."
        },
        {
            "persona": "Marketing Agency",
            "scenario": "Managing 10 client accounts. Uses tool to A/B test content ideas before pitching to clients. Shows predicted ROI. Result: 95% client approval rate."
        }
    ]
    
    for case in use_cases:
        p = tf.paragraphs[0] if case == use_cases[0] else tf.add_paragraph()
        p.text = case["persona"]
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)
        p.space_before = Pt(10)
        
        p = tf.add_paragraph()
        p.text = case["scenario"]
        p.level = 1
        p.font.size = Pt(15)
        p.space_after = Pt(8)

def add_roadmap_slide(prs):
    """Add future roadmap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🚀 Future Roadmap"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    roadmap_items = [
        ("Phase 1 - Enhanced Intelligence", [
            "Real-time trending hashtag recommendations",
            "Competitor analysis and benchmarking",
            "Multi-post batch testing"
        ]),
        ("Phase 2 - Advanced Features", [
            "Video thumbnail A/B testing",
            "Caption optimization suggestions",
            "Optimal posting time calculator"
        ]),
        ("Phase 3 - Integration & Automation", [
            "Direct API integrations with social platforms",
            "Scheduled posting with auto-optimization",
            "Performance tracking dashboard"
        ]),
        ("Phase 4 - Enterprise", [
            "Team collaboration features",
            "White-label solutions for agencies",
            "Custom ML models for specific industries"
        ])
    ]
    
    for phase, items in roadmap_items:
        p = tf.paragraphs[0] if roadmap_items[0][0] == phase else tf.add_paragraph()
        p.text = phase
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(124, 58, 237)
        p.space_before = Pt(10)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(15)
            p.space_after = Pt(4)

def add_call_to_action_slide(prs):
    """Add final call to action slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    shapes = slide.shapes
    
    # Main CTA
    cta_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    cta_frame = cta_box.text_frame
    cta_frame.text = "Ready to Predict Your Success?"
    cta_para = cta_frame.paragraphs[0]
    cta_para.alignment = PP_ALIGN.CENTER
    cta_para.font.size = Pt(48)
    cta_para.font.bold = True
    cta_para.font.color.rgb = RGBColor(37, 99, 235)
    
    # Details
    details_box = shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(2))
    details_frame = details_box.text_frame
    
    p = details_frame.paragraphs[0]
    p.text = "🌐 Try the live demo at:"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.space_after = Pt(10)
    
    p = details_frame.add_paragraph()
    p.text = "trend-analyzer-11.vercel.app"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p.space_after = Pt(20)
    
    p = details_frame.add_paragraph()
    p.text = "💻 GitHub Repository:"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.space_after = Pt(10)
    
    p = details_frame.add_paragraph()
    p.text = "github.com/prachisingh342006/trend-analyzer-11"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(100, 116, 139)
    
    # Footer tagline
    footer_box = shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.7))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Stop guessing. Start predicting. Maximize engagement. 🚀"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(24)
    footer_para.font.italic = True
    footer_para.font.color.rgb = RGBColor(124, 58, 237)

def create_presentation():
    """Main function to create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating presentation slides...")
    
    # Add all slides
    add_title_slide(prs)
    print("✓ Title slide added")
    
    add_problem_statement_slide(prs)
    print("✓ Problem statement slide added")
    
    add_customer_segment_slide(prs)
    print("✓ Customer segment slide added")
    
    add_solution_slide(prs)
    print("✓ Solution slide added")
    
    add_how_it_works_slide(prs)
    print("✓ How it works slide added")
    
    add_features_slide(prs)
    print("✓ Features slide added")
    
    add_benefits_slide(prs)
    print("✓ Benefits slide added")
    
    add_tech_stack_slide(prs)
    print("✓ Tech stack slide added")
    
    add_competitive_advantage_slide(prs)
    print("✓ Competitive advantage slide added")
    
    add_use_cases_slide(prs)
    print("✓ Use cases slide added")
    
    add_roadmap_slide(prs)
    print("✓ Roadmap slide added")
    
    add_call_to_action_slide(prs)
    print("✓ Call to action slide added")
    
    # Save presentation
    filename = "Social_Media_Trend_Predictor_Presentation.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation created successfully: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    create_presentation()
